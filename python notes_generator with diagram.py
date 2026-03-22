# ==========================================
# 🤖 JARVIS LEVEL AI STUDY ASSISTANT
# ==========================================

import streamlit as st
import fitz
from pptx import Presentation
from docx import Document
import matplotlib.pyplot as plt
import re
import os
import numpy as np

from sentence_transformers import SentenceTransformer
from duckduckgo_search import DDGS

from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet

# ==========================================
# 🔐 OPENAI (OPTIONAL BUT POWERFUL)
# ==========================================

USE_AI = False
try:
    from openai import OpenAI
    if os.getenv("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"):
        client = OpenAI()
        USE_AI = True
except:
    pass

# ==========================================
# 🧠 EMBEDDING MODEL
# ==========================================

model = SentenceTransformer('all-MiniLM-L6-v2')

# ==========================================
# 🎨 UI CONFIG
# ==========================================

st.set_page_config(page_title="DEEPAN AI", layout="wide")

st.markdown("""
<style>
.chat-user {background:#1f6feb;padding:10px;border-radius:10px;margin:5px;text-align:right;}
.chat-ai {background:#30363d;padding:10px;border-radius:10px;margin:5px;text-align:left;}
</style>
""", unsafe_allow_html=True)

# ==========================================
# 📄 FILE READERS
# ==========================================

def read_pdf(file):
    text = ""
    pdf = fitz.open(stream=file.read(), filetype="pdf")
    for page in pdf:
        text += page.get_text()
    return text

def read_ppt(file):
    text = ""
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_docx(file):
    return "\n".join([p.text for p in Document(file).paragraphs])

def read_txt(file):
    return file.read().decode("utf-8")

# ==========================================
# 🧠 TEXT PROCESSING
# ==========================================

def split_text(text):
    return [c.strip() for c in text.split(".") if len(c) > 40]

def embed_chunks(chunks):
    return model.encode(chunks)

def semantic_search(question, chunks, embeddings):
    q_emb = model.encode([question])[0]
    scores = np.dot(embeddings, q_emb)
    top_idx = np.argsort(scores)[-5:][::-1]
    return [chunks[i] for i in top_idx]

# ==========================================
# 🌐 WEB SEARCH
# ==========================================

def web_search(query):
    results = []
    with DDGS() as ddgs:
        for r in ddgs.text(query, max_results=3):
            results.append(r["body"])
    return " ".join(results)

# ==========================================
# 🤖 JARVIS CHAT
# ==========================================

def get_answer(question, chunks, embeddings, history):

    doc_context = " ".join(semantic_search(question, chunks, embeddings))
    web_context = web_search(question)

    full_context = doc_context + "\n" + web_context

    if USE_AI:
        try:
            res = client.chat.completions.create(
                model="gpt-5-mini",
                messages=[{
                    "role": "user",
                    "content": f"""
You are JARVIS AI.

- Think deeply
- Answer differently each time
- Be clear and smart

Context:
{full_context}

Question:
{question}
"""
                }]
            )
            return res.choices[0].message.content
        except:
            pass

    return f"""
📌 From Notes:
{doc_context[:400]}

🌐 From Web:
{web_context[:400]}
"""

# ==========================================
# 📅 STUDY PLAN
# ==========================================

def generate_plan(topics, days):
    if not topics:
        return "No topics found"

    per_day = max(1, len(topics)//days)
    plan = []
    idx = 0

    for d in range(days):
        today = topics[idx:idx+per_day]
        idx += per_day

        plan.append(f"Day {d+1}: {', '.join(today)} + Revision")

    return "\n".join(plan)

# ==========================================
# 🗺 MINDMAP
# ==========================================

def create_mindmap(topics):
    if not topics:
        return None

    fig, ax = plt.subplots(figsize=(8,6))

    root = topics[0]
    ax.text(0.5, 0.9, root, ha='center',
            bbox=dict(boxstyle="round", fc="green"))

    spacing = 0.8 / max(1, len(topics)-1)

    for i, t in enumerate(topics[1:]):
        x = 0.1 + i * spacing
        ax.text(x, 0.6, t, ha='center',
                bbox=dict(boxstyle="round", fc="blue"))

        ax.annotate("", xy=(0.5,0.85), xytext=(x,0.65),
                    arrowprops=dict(arrowstyle="->"))

    ax.axis('off')

    path = "mindmap.png"
    plt.savefig(path, bbox_inches='tight')
    plt.close()

    return path

# ==========================================
# 📝 QUIZ
# ==========================================

def generate_quiz(text):
    sents = [s for s in text.split(".") if len(s) > 40]
    return [f"What is: {s[:50]}?" for s in sents[:5]]

# ==========================================
# 📤 EXPORT
# ==========================================

def export_pdf(content):
    doc = SimpleDocTemplate("notes.pdf")
    styles = getSampleStyleSheet()
    doc.build([Paragraph(content, styles["Normal"])])
    return "notes.pdf"

# ==========================================
# 🎛 SIDEBAR
# ==========================================

st.sidebar.title("⚙️ Controls")

file = st.sidebar.file_uploader("Upload File", type=["pdf","pptx","docx","txt"])
days = st.sidebar.slider("Study Days", 1, 30, 3)

# ==========================================
# MAIN UI
# ==========================================

st.title("🤖 JARVIS AI Study Assistant")

if file:

    if file.name.endswith(".pdf"):
        text = read_pdf(file)
    elif file.name.endswith(".pptx"):
        text = read_ppt(file)
    elif file.name.endswith(".docx"):
        text = read_docx(file)
    else:
        text = read_txt(file)

    chunks = split_text(text)

    if "embeddings" not in st.session_state:
        st.session_state.embeddings = embed_chunks(chunks)

    if "chat" not in st.session_state:
        st.session_state.chat = []

    # Preview
    with st.expander("📄 Preview"):
        st.write(text[:1200])

    # Mindmap
    st.subheader("🗺 Mind Map")
    topics = list(set(re.findall(r'\b[A-Z][a-zA-Z]{4,}\b', text)))[:8]
    mind = create_mindmap(topics)
    if mind:
        st.image(mind)

    # Study Plan
    st.subheader("📅 Study Plan")
    plan = generate_plan(topics, days)
    st.write(plan)

    # Quiz
    st.subheader("📝 Quiz")
    for q in generate_quiz(text):
        st.write("❓", q)

    # Chat
    st.subheader("💬 Chat")

    for q, a in st.session_state.chat:
        st.markdown(f"<div class='chat-user'>🧑 {q}</div>", unsafe_allow_html=True)
        st.markdown(f"<div class='chat-ai'>🤖 {a}</div>", unsafe_allow_html=True)

    user_input = st.chat_input("Ask anything...")

    if user_input:
        ans = get_answer(user_input, chunks, st.session_state.embeddings, st.session_state.chat)
        st.session_state.chat.append((user_input, ans))
        st.rerun()

    # Download
    notes = f"Topics:\n{topics}\n\nPlan:\n{plan}"
    pdf = export_pdf(notes)

    with open(pdf, "rb") as f:
        st.download_button("📤 Download Notes", f, "notes.pdf")

else:
    st.info("Upload a file to activate Jarvis 🚀")
